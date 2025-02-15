from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from .models import Event, Participant
import pandas as pd
from .convter import *
from pptx import Presentation
from django.core.mail import send_mail, EmailMessage
import requests
import os

def index(request):
    return render(request, 'index.html')

@login_required
def create(request):
    if request.method == "POST":
        csv = request.FILES.get('csv')
        temp = request.FILES.get('template')
        event_name = request.POST.get('event_name')
        
        event = Event(user = request.user,
            event_name = event_name,
            csv_file = csv,
            template = temp)
        event.save()

        return redirect(f"/certificate/{event.id}/{event.slug}")

    return render(request, 'certificate/create_event.html')

@login_required
def delete_event(request, id, slug):
    event = Event.objects.filter(slug=slug, id=id).first()
    if event.user == request.user:
        event.delete()
    return redirect('view_certificate_status')

@login_required
def track(request, id, slug):
    event = Event.objects.filter(slug=slug, id=id).first()
    if event.message:

        return render(request, 'certificate/track.html', {
            'event_name': event.event_name,
            'event_date': event.date,
            'participat_details': Participant.objects.filter(event=event)
            })

    prs = Presentation(event.template)
    st=""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                st = st + shape.text
                st = st + " "
    li = st.split()
    tags = []
    for i in li:
        if i[0] == "<" and i[-1] == ">":
            tags.append(i)

    if request.method == "POST":
        email_col = request.POST.get('emails')
        subject = request.POST.get('subject')
        mess = request.POST.get('mess')
        values = [(tag, request.POST.get(f'type_{tag}'), request.POST.get(f'input_{tag}')) for tag in tags]    
        
        event.email_column = email_col
        event.message = mess
        event.subject = subject
        event.save()

        df=pd.read_csv(event.csv_file)
        df_len=df.shape
        i=0

        # data = {
        #     "client_id":"YOUR_GOOGLE_API_CLIENT_ID",
        #     "client_secret":"YOUR_CLIENT_SECRET",
        #     "refresh_token": "GOOGLE DRIVE REFRESH TOKEN",
        #     'grant_type': 'refresh_token'
        #     }
        # a = requests.post("https://www.googleapis.com/oauth2/v4/token", data)
        # token = f"Bearer {dict(a.json()).get('access_token')}"
        li=["First","Second","Third"]
        while i < df_len[0]:
            prs = Presentation(event.template)
            j=""
            if i<9:
                j="00"
            elif i>=9 and i < 99 :
                j="0"
            
            for tag, v_type, value in values:
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            if tag in shape.text:
                                text_frame = shape.text_frame
                                
                                # Get replacement text based on type
                                if v_type == 'text':
                                    replacement = value
                                elif v_type == 'date':
                                    replacement = '/'.join(value.split('-')[::-1])
                                elif v_type == 'csv':
                                    replacement = str(df.loc[i, value])
                                elif v_type == "auto":
                                    replacement = f"{value}/{j}{str(i+1)}"
                                else:
                                    continue
                                
                                # Replace text in each paragraph
                                for paragraph in text_frame.paragraphs:
                                    if tag in paragraph.text:
                                        paragraph.text = paragraph.text.replace(tag, replacement)
                                            
            
            s_name = df.loc[i,event.email_column].split('@')[0]
            # prs.save(s_name+".pptx")
            # f_id = ppt2pdf(s_name+".pptx",s_name, token)
            # r = requests.get(f"https://docs.google.com/presentation/d/{f_id}/export/pdf", allow_redirects=True)
            # open(s_name+'.pdf', 'wb').write(r.content)

            pptx_filename = f"/tmp/{s_name}.pptx"
            pdf_filename = f"/tmp/{s_name}.pdf"
            prs.save(pptx_filename)
            convert_pptx_to_pdf(pptx_filename, pdf_filename)
            try:
                mail = EmailMessage(subject,
                    f"Hello, {s_name} \n{mess}",
                    settings.EMAIL_HOST_USER,
                    [df.loc[i,event.email_column]])
                mail.attach_file(pdf_filename)
                mail.send()
                Participant(event=event, email=df.loc[i,event.email_column], status=True).save()
                os.remove(pdf_filename)
                os.remove(pptx_filename)
            except:
                Participant(event=event, email=df.loc[i,event.email_column], status=False).save()
                os.remove(pdf_filename)
                os.remove(pptx_filename)
            i=i+1

        messages.success(request, "Certificates Sent Successfuly !!")
        return redirect(f"/certificate/{event.id}/{event.slug}")


    return render(request, 'certificate/map_tags_of_template.html',{
        'columns': list(pd.read_csv(event.csv_file).columns),
        'tags': tags,
        })


@login_required
def view_certificate_status(request):
    return render(request, 'certificate/view_certificate_status.html', {
        'events': Event.objects.filter(user=request.user).order_by('-date')
    })
